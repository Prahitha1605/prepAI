from flask import Flask, render_template, request, redirect, url_for, session, flash, jsonify, send_file, send_from_directory, after_this_request
from flask_cors import CORS
from flask_pymongo import PyMongo
from werkzeug.security import generate_password_hash, check_password_hash
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from textblob import TextBlob
import nltk
import textstat
import google.generativeai as genai
from datetime import timedelta
import os
import uuid
import threading
import time
from dotenv import load_dotenv
from feedback_analysis import analyze_responses
from gemini_api import generate_questions
from transcribe import transcribe_video

nltk.download('punkt')

# Flask app setup
app = Flask(__name__)
CORS(app)
app.secret_key = os.urandom(24).hex()
app.permanent_session_lifetime = timedelta(days=1)

# MongoDB Configuration
app.config["MONGO_URI"] = "mongodb://localhost:27017/jics_db"
mongo = PyMongo(app)

# Video upload folder
UPLOAD_FOLDER = "static/videos"
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

# Load environment variables
load_dotenv()
GEMINI_API_KEY = os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=GEMINI_API_KEY)

# Gemini AI configuration
generation_config = {
    "temperature": 0.7,
    "top_p": 0.9,
    "top_k": 20,
    "max_output_tokens": 100,
    "response_mime_type": "text/plain",
}
model = genai.GenerativeModel(
    model_name="gemini-1.5-flash",
    generation_config=generation_config,
    system_instruction=(
        "You are a resume builder chatbot. Respond concisely, using plain text with line breaks. "
        "Organize information using colons (:) and avoid markdown symbols like **. "
        "Focus on resume suggestions, skills, improvements, and relevant courses."
    )
)

# Initialize MongoDB
def init_db():
    users_collection = mongo.db.users
    users_collection.create_index("username", unique=True)

# Routes
@app.route('/')
def home():
    return render_template('landing_page.html', username=session.get('username'))

@app.route('/register', methods=['GET', 'POST'])
def register():
    if request.method == 'POST':
        users_collection = mongo.db.users
        username = request.form['username']
        email = request.form['email']
        password = request.form['password']
        confirm_password = request.form['confirm_password']

        if password != confirm_password:
            flash('Passwords do not match!', 'error')
            return redirect(url_for('register'))

        existing_user = users_collection.find_one({"username": username})
        if existing_user:
            flash('Username already exists!', 'error')
            return redirect(url_for('register'))

        hashed_password = generate_password_hash(password)
        users_collection.insert_one({"username": username, "email": email, "password": hashed_password})

        flash('Registration successful! Please log in.', 'success')
        return redirect(url_for('login'))

    return render_template('register.html')

@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        users_collection = mongo.db.users
        username = request.form['username']
        password = request.form['password']

        user = users_collection.find_one({"username": username})

        if user and check_password_hash(user["password"], password):
            session.permanent = True
            session['username'] = username
            flash('Login successful!', 'success')
            return redirect(url_for('dashboard'))
        else:
            flash('Invalid credentials! Please register if you do not have an account.', 'error')
            return redirect(url_for('login'))

    return render_template('login.html')

@app.route('/landing_page')
def landing_page():
    return render_template('landing_page.html', username=session.get('username'))

@app.route('/dashboard')
def dashboard():
    if 'username' in session:
        return render_template('dashboard.html', username=session['username'])
    flash('You need to log in first!', 'error')
    return redirect(url_for('login'))

@app.route('/resume_intro')
def resume_intro():
    return render_template('resume/resume_intro.html')

@app.route('/create_resume')
def create_resume():
    return render_template('resume/create-resume.html')

@app.route("/enhance-resume", methods=["POST"])
def enhance_resume():
    try:
        data = request.json
        enhanced_sections = {}

        if data.get('education'):
            education_response = model.generate_content(f"Professionally format this education entry:\n{data['education'][0]}")
            enhanced_sections['Education'] = education_response.text.strip()

        if data.get('experience'):
            experience_response = model.generate_content(f"Enhance this internship experience:\n{data['experience'][0]}")
            enhanced_sections['Experience'] = experience_response.text.strip()

        if data.get('skills'):
            skills_response = model.generate_content(f"Refine these skills:\n{', '.join(data['skills'])}")
            enhanced_sections['Skills'] = skills_response.text.strip()

        if data.get('certifications'):
            cert_response = model.generate_content(f"Enhance these certifications:\n{', '.join(data['certifications'])}")
            enhanced_sections['Certifications'] = cert_response.text.strip()

        enhanced_resume_html = ""
        for title, content in enhanced_sections.items():
            enhanced_resume_html += f"<div class='resume-section'><h3>{title}</h3><ul><li>{content}</li></ul></div>"

        return jsonify({"enhanced_resume": enhanced_resume_html or "No enhancements found."})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/chat', methods=['POST'])
def chat():
    user_input = request.json.get('message')
    if not user_input:
        return jsonify({'error': 'No message provided'}), 400
    try:
        response = model.generate_content(f"User said: {user_input}")
        formatted_response = format_bot_response(response.text if response else "I didn't understand that.")
        return jsonify({'response': formatted_response})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

def format_bot_response(text):
    sections = ["Programming Languages", "Frameworks/Libraries", "Databases", "Tools", "Other Skills", "Resume Improvement Suggestions", "Certifications", "Experience", "Education", "Soft Skills", "Projects", "Courses"]
    for section in sections:
        text = text.replace(f"{section}:", f"\n\n*{section}:*\n")
    return text.replace(", ", "\n- ").strip()

@app.route('/profile')
def profile():
    if 'username' not in session:
        return jsonify({"error": "Not logged in"}), 403
    users_collection = mongo.db.users
    user = users_collection.find_one({"username": session['username']})
    if user:
        return jsonify({"id": str(user["_id"]), "username": user["username"], "email": user["email"]})
    return jsonify({"error": "User not found"}), 404

@app.route("/logout")
def logout():
    session.pop("username", None)
    flash("You have been logged out successfully!", "success")
    return redirect(url_for("home"))

@app.route("/evaluate_answer", methods=["POST"])
def evaluate_answer():
    user_answer = request.json.get("answer")
    if not user_answer:
        return jsonify({"feedback": "Please enter a valid answer."})
    blob = TextBlob(user_answer)
    polarity = blob.sentiment.polarity
    if polarity > 0.3:
        feedback = "Great positive tone! Try to add specific examples."
    elif polarity < -0.1:
        feedback = "The answer sounds a bit negative. Try to be more optimistic."
    else:
        feedback = "Good start! Try to include more detail or enthusiasm."
    return jsonify({"feedback": feedback})

@app.route("/presentation_guidance", methods=["GET", "POST"])
def presentation_guidance():
    def apply_template(slide, template_name):
        title_shape = slide.shapes.title
        if title_shape:
            title_text_frame = title_shape.text_frame
            title_text_frame.paragraphs[0].font.size = Pt(36)
            title_text_frame.paragraphs[0].font.bold = True
            title_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.size = Pt(20)
                    paragraph.space_after = Pt(5)
        if template_name == "professional":
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(220, 230, 241)
        elif template_name == "dark":
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(34, 34, 34)
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        else:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    def generate_guidance(text):
        blob = TextBlob(text)
        sentiment = blob.sentiment.polarity
        length = len(text.split())
        guidance = []
        readability = textstat.flesch_reading_ease(text)
        if readability < 50:
            guidance.append("Difficult to read. Simplify language.")
        elif readability > 70:
            guidance.append("Easy to read. Great job!")
        if length > 20:
            guidance.append("Use visuals to support long text.")
        if len(blob.sentences) > 3:
            guidance.append("Avoid too much text on one slide.")
        if length > 30:
            guidance.append("Break content into bullet points.")
        elif length < 5:
            guidance.append("Content is short. Add more details.")
        if sentiment < -0.2:
            guidance.append("Tone is negative. Try to rephrase.")
        elif sentiment > 0.2:
            guidance.append("Tone is positive. Keep it up!")
        else:
            guidance.append("Tone is neutral. Add more emotion.")
        return guidance if guidance else ["Looks good!"]

    def generate_ppt(title, slides_content, template):
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        apply_template(slide, template)
        guidance_summary = []
        for slide_title, slide_text in slides_content.items():
            slide = prs.slides.add_slide(content_slide_layout)
            slide.shapes.title.text = slide_title
            slide.shapes.placeholders[1].text = slide_text
            apply_template(slide, template)
            guidance = generate_guidance(slide_text)
            guidance_summary.append(f"{slide_title}: {', '.join(guidance)}")
        summary_slide = prs.slides.add_slide(content_slide_layout)
        summary_slide.shapes.title.text = "Guidance Summary"
        textbox = summary_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(5))
        text_frame = textbox.text_frame
        text_frame.text = "\n".join(guidance_summary)
        ppt_file = f"presentation_{uuid.uuid4().hex}.pptx"
        prs.save(ppt_file)
        return ppt_file

    if request.method == "POST":
        title = request.form.get("title")
        template = request.form.get("template")
        slides_content = {}
        for i in range(1, 6):
            slide_title = request.form.get(f"slide_title_{i}")
            slide_text = request.form.get(f"slide_text_{i}")
            if slide_title and slide_text:
                slides_content[slide_title] = slide_text
        ppt_file = generate_ppt(title, slides_content, template)

        @after_this_request
        def remove_file(response):
            def delayed_delete(file_path):
                time.sleep(5)
                try:
                    os.remove(file_path)
                except Exception as e:
                    print(f"Error deleting file: {e}")
            threading.Thread(target=delayed_delete, args=(ppt_file,)).start()
            return response

        return send_file(ppt_file, as_attachment=True)

    return render_template("p_index.html")

@app.route("/index")
def index():
    return render_template("index.html")

@app.route("/interview")
def interview_page():
    return render_template("interview.html")

@app.route("/feedback")
def feedback_page():
    return render_template("feedback.html")

@app.route("/generate_questions", methods=["POST"])
def generate_interview_questions():
    try:
        data = request.json
        career_data = data.get("career_info")
        if not career_data:
            return jsonify({"error": "No career_info provided"}), 400
        questions = generate_questions(career_data)
        return jsonify({"questions": questions})
    except Exception as e:
        print(f"Error in generate_questions: {str(e)}")
        return jsonify({"error": f"Failed to generate questions: {str(e)}"}), 500

@app.route("/uploadvideo", methods=["POST"])
def uploadvideo():
    if "video" not in request.files:
        print("No video file provided in request")
        return jsonify({"error": "No video file provided"}), 400

    file = request.files["video"]
    if file.filename == "":
        print("Empty filename in video upload")
        return jsonify({"error": "No video file selected"}), 400

    try:
        file_path = os.path.join(app.config["UPLOAD_FOLDER"], file.filename)
        file.save(file_path)
        if not os.path.exists(file_path):
            print(f"Failed to save file at {file_path}")
            return jsonify({"error": "Failed to save video file"}), 500
        relative_file_path = file_path.replace(os.path.abspath("."), "").replace("\\", "/").lstrip("/")
        print(f"Video uploaded successfully: {relative_file_path}")
        return jsonify({"message": "File uploaded successfully!", "file_path": relative_file_path})
    except Exception as e:
        print(f"Error saving video: {str(e)}")
        return jsonify({"error": f"Failed to save video: {str(e)}"}), 500

@app.route("/analyze_responses", methods=["POST"])
def analyze_video():
    start_time = time.time()
    try:
        data = request.json
        file_paths = data.get("file_paths", [])
        if not file_paths:
            print("No file paths provided for analysis")
            return jsonify({"error": "No video files provided for analysis"}), 400

        responses = []
        for idx, file_path in enumerate(file_paths, 1):
            absolute_file_path = os.path.abspath(file_path)
            print(f"Processing file {idx}: {absolute_file_path}")
            if not os.path.exists(absolute_file_path):
                print(f"File not found: {absolute_file_path}")
                responses.append({"text": f"[File not found: {absolute_file_path}]", "error": True})
            else:
                try:
                    file_start_time = time.time()
                    response_text = transcribe_video(absolute_file_path)
                    file_end_time = time.time()
                    print(f"Transcription for {absolute_file_path} took {file_end_time - file_start_time:.2f} seconds: {response_text}")
                    responses.append({"text": response_text, "error": False})
                except Exception as e:
                    print(f"Error transcribing {absolute_file_path}: {str(e)}")
                    responses.append({"text": f"[Transcription Error: {str(e)}]", "error": True})

        try:
            analysis_start_time = time.time()
            valid_responses = [r["text"] for r in responses if not r["error"]]
            analysis_results = analyze_responses(valid_responses)
            analysis_end_time = time.time()
            print(f"Analysis took {analysis_end_time - analysis_start_time:.2f} seconds")

            final_results = []
            for idx, (resp, analysis) in enumerate(zip(responses, analysis_results + [{}] * (len(responses) - len(analysis_results))), 1):
                if resp["error"]:
                    final_results.append({
                        "question_number": idx,
                        "response_text": resp["text"],
                        "score": 0,
                        "feedback": "Unable to analyze due to transcription error"
                    })
                else:
                    analysis["question_number"] = idx
                    final_results.append(analysis)
            end_time = time.time()
            print(f"Total /analyze_responses took {end_time - start_time:.2f} seconds")
            print(f"Final results: {final_results}")
            return jsonify({"status": "success", "results": final_results})
        except Exception as e:
            print(f"Analysis failed: {str(e)}")
            return jsonify({"status": "error", "message": f"Analysis failed: {str(e)}"}), 500
    except Exception as e:
        print(f"Error in analyze_responses: {str(e)}")
        return jsonify({"status": "error", "message": f"Request processing failed: {str(e)}"}), 500

@app.route("/test_transcription", methods=["POST"])
def test_transcription():
    try:
        data = request.json
        file_path = data.get("file_path")
        if not file_path:
            print("No file path provided for test transcription")
            return jsonify({"error": "No file path provided"}), 400
        absolute_file_path = os.path.abspath(file_path)
        if not os.path.exists(absolute_file_path):
            print(f"File not found for test transcription: {absolute_file_path}")
            return jsonify({"error": f"File not found: {absolute_file_path}"}), 404
        start_time = time.time()
        transcription = transcribe_video(absolute_file_path)
        end_time = time.time()
        print(f"Test transcription took {end_time - start_time:.2f} seconds: {transcription}")
        return jsonify({"transcription": transcription})
    except Exception as e:
        print(f"Error in test_transcription: {str(e)}")
        return jsonify({"error": f"Transcription failed: {str(e)}"}), 500

# Cache control
@app.before_request
def no_cache():
    if 'username' not in session and request.endpoint in ['dashboard', 'profile']:
        flash("You need to log in first!", "error")
        return redirect(url_for('login'))

@app.after_request
def add_header(response):
    response.headers["Cache-Control"] = "no-store, no-cache, must-revalidate, max-age=0"
    response.headers["Pragma"] = "no-cache"
    response.headers["Expires"] = "0"
    return response

if __name__ == '__main__':
    init_db()
    app.run(debug=True, host='0.0.0.0', port=8080)